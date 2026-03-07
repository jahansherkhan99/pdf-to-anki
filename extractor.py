"""
extractor.py — Extract raw text from PDF and PPTX files.

Returns a list of (page_or_slide_number, text) tuples so the rest of the
pipeline knows where content came from for progress reporting.
"""

from __future__ import annotations

from pathlib import Path
from typing import List, Tuple

import fitz  # PyMuPDF
from pptx import Presentation

from config import MAX_CHUNK_WORDS

# ---------------------------------------------------------------------------
# Public types
# ---------------------------------------------------------------------------

Page = Tuple[int, str]   # (1-based number, extracted text)
Chunk = dict             # {"start": int, "end": int, "text": str}


# ---------------------------------------------------------------------------
# Extraction helpers
# ---------------------------------------------------------------------------

def _extract_pdf(path: str) -> List[Page]:
    """Extract text from every page of a PDF, skipping blank pages."""
    doc = fitz.open(path)
    pages: List[Page] = []
    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        if text:
            pages.append((page_num, text))
    doc.close()
    return pages


def _extract_pptx(path: str) -> List[Page]:
    """Extract text from every slide of a PPTX, skipping blank slides."""
    prs = Presentation(path)
    slides: List[Page] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    parts.append(line)
        text = "\n".join(parts)
        if text:
            slides.append((slide_num, text))
    return slides


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def extract(path: str) -> List[Page]:
    """
    Extract text from a PDF or PPTX file.

    Returns a list of (number, text) tuples where *number* is the 1-based
    page (PDF) or slide (PPTX) index.

    Raises ValueError for unsupported file types.
    """
    suffix = Path(path).suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix in (".pptx", ".ppt"):
        return _extract_pptx(path)
    raise ValueError(
        f"Unsupported file type '{suffix}'. Only .pdf and .pptx are accepted."
    )


def chunk_pages(pages: List[Page], max_words: int = MAX_CHUNK_WORDS) -> List[Chunk]:
    """
    Group pages/slides into chunks whose combined word count does not exceed
    *max_words*.  A single page that is already larger than *max_words* is
    kept as its own chunk (we never truncate source content).

    Each chunk is a dict:
        {"start": <first page num>, "end": <last page num>, "text": <combined text>}
    """
    chunks: List[Chunk] = []
    current_pages: List[str] = []
    current_words = 0
    current_start: int | None = None

    for page_num, text in pages:
        word_count = len(text.split())

        # Flush current chunk when adding this page would exceed the limit
        # (unless the chunk is empty — we always include at least one page).
        if current_pages and current_words + word_count > max_words:
            chunks.append(
                {
                    "start": current_start,
                    "end": page_num - 1,
                    "text": "\n\n".join(current_pages),
                }
            )
            current_pages = []
            current_words = 0
            current_start = None

        if not current_pages:
            current_start = page_num
        current_pages.append(text)
        current_words += word_count

    # Flush remaining pages
    if current_pages:
        chunks.append(
            {
                "start": current_start,
                "end": pages[-1][0],
                "text": "\n\n".join(current_pages),
            }
        )

    return chunks
